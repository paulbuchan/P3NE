import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.exc import PackageNotFoundError

class PPTXNotesExtractor:
    def __init__(self, root):
        self.root = root
        self.root.withdraw()  # Hide the main window
        self.open_file_dialog()

    def extract_notes(self, ppt_path):
        try:
            prs = Presentation(ppt_path)
            notes = []
            for i, slide in enumerate(prs.slides, start=1):
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else ''
                    notes.append(f"[Slide {i}]\n{text}\n")
            return notes
        except PackageNotFoundError:
            messagebox.showerror("Error", "Invalid PowerPoint file.")
            return None
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {e}")
            return None

    def save_notes(self, notes):
        if not notes:
            messagebox.showinfo("No Notes", "No notes were found in the presentation.")
            return False

        file_path = filedialog.asksaveasfilename(defaultextension=".txt",
                                                filetypes=[("Text files", "*.txt")])

        # Check if the save dialog was cancelled
        if not file_path:
            messagebox.showinfo("Cancelled", "Operation cancelled by user.")
            return False

        with open(file_path, 'w', encoding='utf-8') as file:
            file.writelines(notes)
        messagebox.showinfo("Success", "Notes saved successfully!")
        return True



    def open_file_dialog(self):
        file_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
        if file_path:
            notes = self.extract_notes(file_path)
            if notes is not None:
                saved = self.save_notes(notes)
                self.root.destroy()  # Close the app after attempting to save
        else:
            self.root.destroy()  # Close the app if file dialog is cancelled



root = tk.Tk()
app = PPTXNotesExtractor(root)
root.mainloop()
